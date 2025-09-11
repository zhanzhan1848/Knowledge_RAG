# -*- coding: utf-8 -*-
"""
文档预览服务模块

提供文档预览功能，支持多种文件格式的预览生成
包括PDF、图片、文本文件等的预览处理
"""

import os
import io
import tempfile
from typing import Optional, Tuple, Dict, Any
from pathlib import Path
import logging

from PIL import Image, ImageDraw, ImageFont
import fitz  # PyMuPDF
from docx import Document as DocxDocument
import openpyxl
from pptx import Presentation
import markdown
from fastapi import HTTPException
from fastapi.responses import StreamingResponse

from .config import get_settings
from .storage import get_minio_storage

logger = logging.getLogger(__name__)


class DocumentPreviewService:
    """文档预览服务类
    
    提供各种文档格式的预览功能
    """
    
    def __init__(self):
        self.settings = get_settings()
        self.storage = get_minio_storage()
        self.supported_formats = {
            'application/pdf': self._preview_pdf,
            'image/jpeg': self._preview_image,
            'image/png': self._preview_image,
            'image/gif': self._preview_image,
            'image/webp': self._preview_image,
            'text/plain': self._preview_text,
            'text/markdown': self._preview_markdown,
            'text/html': self._preview_html,
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': self._preview_docx,
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': self._preview_xlsx,
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': self._preview_pptx,
            'application/json': self._preview_json,
            'text/csv': self._preview_csv
        }
    
    def is_previewable(self, content_type: str) -> bool:
        """检查文件类型是否支持预览
        
        Args:
            content_type: 文件MIME类型
            
        Returns:
            bool: 是否支持预览
        """
        return content_type in self.supported_formats
    
    async def generate_preview(self, 
                             file_path: str, 
                             content_type: str,
                             page: int = 1,
                             size: Tuple[int, int] = (800, 600)) -> StreamingResponse:
        """生成文档预览
        
        Args:
            file_path: MinIO中的文件路径
            content_type: 文件MIME类型
            page: 页码（对于多页文档）
            size: 预览图片尺寸
            
        Returns:
            StreamingResponse: 预览图片流
        """
        if not self.is_previewable(content_type):
            raise HTTPException(
                status_code=400, 
                detail=f"不支持的文件类型: {content_type}"
            )
        
        try:
            # 从MinIO下载文件到临时目录
            with tempfile.NamedTemporaryFile() as temp_file:
                file_data = await self.storage.download_file(file_path)
                temp_file.write(file_data)
                temp_file.flush()
                
                # 生成预览
                preview_func = self.supported_formats[content_type]
                preview_image = await preview_func(temp_file.name, page, size)
                
                # 转换为字节流
                img_buffer = io.BytesIO()
                preview_image.save(img_buffer, format='PNG')
                img_buffer.seek(0)
                
                return StreamingResponse(
                    io.BytesIO(img_buffer.getvalue()),
                    media_type="image/png",
                    headers={"Cache-Control": "public, max-age=3600"}
                )
                
        except Exception as e:
            logger.error(f"生成预览失败: {str(e)}")
            raise HTTPException(
                status_code=500,
                detail=f"生成预览失败: {str(e)}"
            )
    
    async def _preview_pdf(self, file_path: str, page: int, size: Tuple[int, int]) -> Image.Image:
        """预览PDF文件
        
        Args:
            file_path: 文件路径
            page: 页码
            size: 图片尺寸
            
        Returns:
            Image.Image: 预览图片
        """
        doc = fitz.open(file_path)
        
        if page > len(doc):
            page = 1
        
        page_obj = doc[page - 1]
        
        # 计算缩放比例
        mat = fitz.Matrix(2.0, 2.0)  # 2倍缩放
        pix = page_obj.get_pixmap(matrix=mat)
        
        # 转换为PIL Image
        img_data = pix.tobytes("ppm")
        img = Image.open(io.BytesIO(img_data))
        
        # 调整尺寸
        img.thumbnail(size, Image.Resampling.LANCZOS)
        
        doc.close()
        return img
    
    async def _preview_image(self, file_path: str, page: int, size: Tuple[int, int]) -> Image.Image:
        """预览图片文件
        
        Args:
            file_path: 文件路径
            page: 页码（图片忽略）
            size: 图片尺寸
            
        Returns:
            Image.Image: 预览图片
        """
        img = Image.open(file_path)
        
        # 转换为RGB模式（如果需要）
        if img.mode != 'RGB':
            img = img.convert('RGB')
        
        # 调整尺寸
        img.thumbnail(size, Image.Resampling.LANCZOS)
        
        return img
    
    async def _preview_text(self, file_path: str, page: int, size: Tuple[int, int]) -> Image.Image:
        """预览文本文件
        
        Args:
            file_path: 文件路径
            page: 页码
            size: 图片尺寸
            
        Returns:
            Image.Image: 预览图片
        """
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        return self._create_text_image(content, size)
    
    async def _preview_markdown(self, file_path: str, page: int, size: Tuple[int, int]) -> Image.Image:
        """预览Markdown文件
        
        Args:
            file_path: 文件路径
            page: 页码
            size: 图片尺寸
            
        Returns:
            Image.Image: 预览图片
        """
        with open(file_path, 'r', encoding='utf-8') as f:
            md_content = f.read()
        
        # 转换为HTML然后提取文本
        html = markdown.markdown(md_content)
        # 简单的HTML标签移除
        import re
        text = re.sub(r'<[^>]+>', '', html)
        
        return self._create_text_image(text, size)
    
    async def _preview_html(self, file_path: str, page: int, size: Tuple[int, int]) -> Image.Image:
        """预览HTML文件
        
        Args:
            file_path: 文件路径
            page: 页码
            size: 图片尺寸
            
        Returns:
            Image.Image: 预览图片
        """
        with open(file_path, 'r', encoding='utf-8') as f:
            html_content = f.read()
        
        # 简单的HTML标签移除
        import re
        text = re.sub(r'<[^>]+>', '', html_content)
        
        return self._create_text_image(text, size)
    
    async def _preview_docx(self, file_path: str, page: int, size: Tuple[int, int]) -> Image.Image:
        """预览Word文档
        
        Args:
            file_path: 文件路径
            page: 页码
            size: 图片尺寸
            
        Returns:
            Image.Image: 预览图片
        """
        doc = DocxDocument(file_path)
        
        # 提取文本内容
        text_content = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content.append(paragraph.text)
        
        content = '\n'.join(text_content)
        return self._create_text_image(content, size)
    
    async def _preview_xlsx(self, file_path: str, page: int, size: Tuple[int, int]) -> Image.Image:
        """预览Excel文件
        
        Args:
            file_path: 文件路径
            page: 页码（工作表索引）
            size: 图片尺寸
            
        Returns:
            Image.Image: 预览图片
        """
        workbook = openpyxl.load_workbook(file_path)
        
        # 获取工作表
        sheet_names = workbook.sheetnames
        if page > len(sheet_names):
            page = 1
        
        sheet = workbook[sheet_names[page - 1]]
        
        # 提取数据
        content_lines = []
        for row in sheet.iter_rows(max_row=20, max_col=10, values_only=True):
            row_data = [str(cell) if cell is not None else '' for cell in row]
            if any(row_data):  # 跳过空行
                content_lines.append('\t'.join(row_data))
        
        content = '\n'.join(content_lines)
        return self._create_text_image(content, size)
    
    async def _preview_pptx(self, file_path: str, page: int, size: Tuple[int, int]) -> Image.Image:
        """预览PowerPoint文件
        
        Args:
            file_path: 文件路径
            page: 页码（幻灯片索引）
            size: 图片尺寸
            
        Returns:
            Image.Image: 预览图片
        """
        prs = Presentation(file_path)
        
        if page > len(prs.slides):
            page = 1
        
        slide = prs.slides[page - 1]
        
        # 提取幻灯片文本
        text_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if shape.text.strip():
                    text_content.append(shape.text)
        
        content = '\n'.join(text_content)
        return self._create_text_image(content, size)
    
    async def _preview_json(self, file_path: str, page: int, size: Tuple[int, int]) -> Image.Image:
        """预览JSON文件
        
        Args:
            file_path: 文件路径
            page: 页码
            size: 图片尺寸
            
        Returns:
            Image.Image: 预览图片
        """
        import json
        
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        # 格式化JSON
        content = json.dumps(data, indent=2, ensure_ascii=False)
        return self._create_text_image(content, size)
    
    async def _preview_csv(self, file_path: str, page: int, size: Tuple[int, int]) -> Image.Image:
        """预览CSV文件
        
        Args:
            file_path: 文件路径
            page: 页码
            size: 图片尺寸
            
        Returns:
            Image.Image: 预览图片
        """
        import csv
        
        content_lines = []
        with open(file_path, 'r', encoding='utf-8') as f:
            reader = csv.reader(f)
            for i, row in enumerate(reader):
                if i >= 20:  # 只显示前20行
                    break
                content_lines.append('\t'.join(row))
        
        content = '\n'.join(content_lines)
        return self._create_text_image(content, size)
    
    def _create_text_image(self, text: str, size: Tuple[int, int]) -> Image.Image:
        """创建文本预览图片
        
        Args:
            text: 文本内容
            size: 图片尺寸
            
        Returns:
            Image.Image: 文本图片
        """
        # 创建白色背景图片
        img = Image.new('RGB', size, color='white')
        draw = ImageDraw.Draw(img)
        
        # 尝试加载字体
        try:
            # macOS系统字体
            font = ImageFont.truetype('/System/Library/Fonts/PingFang.ttc', 14)
        except:
            try:
                # 备用字体
                font = ImageFont.truetype('/System/Library/Fonts/Arial.ttf', 14)
            except:
                # 默认字体
                font = ImageFont.load_default()
        
        # 文本换行处理
        lines = []
        words = text.split('\n')
        
        for line in words[:30]:  # 最多显示30行
            if len(line) > 80:  # 每行最多80个字符
                line = line[:77] + '...'
            lines.append(line)
        
        # 绘制文本
        y_offset = 10
        line_height = 20
        
        for line in lines:
            if y_offset + line_height > size[1] - 10:
                break
            
            draw.text((10, y_offset), line, fill='black', font=font)
            y_offset += line_height
        
        return img


# 全局预览服务实例
_preview_service = None


def get_preview_service() -> DocumentPreviewService:
    """获取文档预览服务实例
    
    Returns:
        DocumentPreviewService: 预览服务实例
    """
    global _preview_service
    if _preview_service is None:
        _preview_service = DocumentPreviewService()
    return _preview_service